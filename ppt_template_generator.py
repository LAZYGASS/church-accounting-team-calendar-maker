"""
재정집행일정 PPT 자동화 템플릿
- 원본 PPT의 디자인 서식을 추출하여 자동 생성
- 규정: 둘째주/넷째주 화요일 예산지급
- 결재일: 지급 전주 금요일/토요일
- 운영위원회의: 마지막주 전주 일요일
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
import calendar
from datetime import datetime, timedelta
import os

class BudgetCalendarTemplate:
    """예산집행 캘린더 PPT 템플릿"""
    
    # 색상 팔레트 (튜플로 저장, 사용 시 RGBColor로 변환)
    _COLOR_BLUE_RGB = (0, 97, 255)      # #0061FF - 집행일, 월 숫자
    _COLOR_PINK_RGB = (253, 121, 168)   # #FD79A8 - 결재일
    _COLOR_SUNDAY_RGB = (253, 121, 168) # 일요일 색상
    _COLOR_SATURDAY_RGB = (68, 114, 196) # #4472C4 - 토요일 파란색
    _COLOR_WHITE_RGB = (255, 255, 255)
    _COLOR_BLACK_RGB = (0, 0, 0)
    _COLOR_BG_RGB = (244, 248, 251)     # #F4F8FB - 배경색
    _COLOR_PREV_MONTH_RGB = (192, 211, 222) # #C0D3DE - 지난달 날짜
    _COLOR_GOLD_RGB = (255, 215, 0)     # 황금색 - 운영위원회

    # RGBColor 객체
    COLOR_BLUE = RGBColor(*_COLOR_BLUE_RGB)
    COLOR_PINK = RGBColor(*_COLOR_PINK_RGB)
    COLOR_SUNDAY = RGBColor(*_COLOR_SUNDAY_RGB)
    COLOR_SATURDAY = RGBColor(*_COLOR_SATURDAY_RGB)
    COLOR_WHITE = RGBColor(*_COLOR_WHITE_RGB)
    COLOR_BLACK = RGBColor(*_COLOR_BLACK_RGB)
    COLOR_BG = RGBColor(*_COLOR_BG_RGB)
    COLOR_PREV_MONTH = RGBColor(*_COLOR_PREV_MONTH_RGB)
    COLOR_GOLD = RGBColor(*_COLOR_GOLD_RGB)
    
    # 레이아웃 설정
    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(7.5)
    
    # 표 설정
    TABLE_LEFT = Inches(0.65)
    TABLE_TOP = Inches(1.59)
    TABLE_WIDTH = Inches(8.71)
    TABLE_HEIGHT = Inches(5.43)
    
    # 폰트 설정
    FONT_TITLE = "NanumSquare Bold"
    FONT_BODY = "NanumSquare Bold"
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
    
    @staticmethod
    def calculate_schedule(year, month):
        """
        규정에 따른 일정 자동 계산

        Returns:
            dict: {
                'execution_days': [둘째주 화요일, 넷째주 화요일],
                'approval_days': [각 집행일 전주 금요일, 토요일],
                'committee_day': 마지막주 전주 일요일
            }
        """
        cal = calendar.monthcalendar(year, month)

        # 1. 예산지급일: 둘째주/넷째주 화요일
        execution_days = []
        tuesday_weeks = []

        for week_idx, week in enumerate(cal):
            tuesday = week[1]  # 화요일 (인덱스 1, Monday=0 기준)
            if tuesday != 0:
                tuesday_weeks.append((week_idx, tuesday))

        # 둘째주, 넷째주 선택
        if len(tuesday_weeks) >= 2:
            execution_days.append(tuesday_weeks[1][1])  # 둘째주
        if len(tuesday_weeks) >= 4:
            execution_days.append(tuesday_weeks[3][1])  # 넷째주

        # 2. 결재일: 각 집행일 전주 금요일, 토요일
        approval_days = []
        for exec_day in execution_days:
            # 집행일로부터 4일 전 (화요일 -> 금요일)
            friday_date = datetime(year, month, exec_day) - timedelta(days=4)
            if friday_date.month == month:
                approval_days.append(friday_date.day)
            # 집행일로부터 3일 전 (화요일 -> 토요일)
            saturday_date = datetime(year, month, exec_day) - timedelta(days=3)
            if saturday_date.month == month:
                approval_days.append(saturday_date.day)
        
        # 3. 운영위원회의: 마지막주 전주 일요일
        committee_day = None
        
        # 마지막 주 찾기
        last_week_idx = len(cal) - 1
        while last_week_idx >= 0 and all(day == 0 for day in cal[last_week_idx]):
            last_week_idx -= 1
        
        # 마지막주 전주 일요일
        if last_week_idx >= 1:
            prev_week = cal[last_week_idx - 1]
            sunday = prev_week[0]  # 일요일 (인덱스 0)
            if sunday != 0:
                committee_day = sunday
        
        return {
            'execution_days': execution_days,
            'approval_days': approval_days,
            'committee_day': committee_day
        }
        
    def create_calendar_slide(self, year, month, 
                             approval_days=None, 
                             execution_days=None,
                             committee_day=None,
                             auto_calculate=True):
        """
        캘린더 슬라이드 생성
        
        Args:
            year: 연도
            month: 월 (1-12)
            approval_days: 결재일 리스트 (None이면 자동 계산)
            execution_days: 집행일 리스트 (None이면 자동 계산)
            committee_day: 운영위원회의 날짜 (None이면 자동 계산)
            auto_calculate: True이면 규정에 따라 자동 계산
        """
        # 자동 계산
        if auto_calculate and (approval_days is None or execution_days is None):
            schedule = self.calculate_schedule(year, month)
            if approval_days is None:
                approval_days = schedule['approval_days']
            if execution_days is None:
                execution_days = schedule['execution_days']
            if committee_day is None:
                committee_day = schedule['committee_day']
        
        # 빈 슬라이드 추가
        blank_layout = self.prs.slide_layouts[6]  # 빈 레이아웃
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 배경색 설정 (#F4F8FB)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_BG
        
        # 1. 월 표시 (왼쪽 상단 큰 숫자)
        self._add_month_number(slide, month)
        
        # 2. 헤더 박스 (예산집행캘린더)
        self._add_header_box(slide, year, month)
        
        # 3. 캘린더 표 생성
        execution_positions, approval_positions = self._add_calendar_table(slide, year, month, approval_days, execution_days, committee_day)

        # 4. 범례 추가 (집행일, 결재일 날짜 밑에 배치)
        self._add_legend(slide, approval_days, execution_days, execution_positions, approval_positions)
        
        return slide
    
    def _add_month_number(self, slide, month):
        """월 숫자 추가 (왼쪽 상단)"""
        left = Inches(0.59)
        top = Inches(0.32)
        width = Inches(1.82)
        height = Inches(1.31)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = str(month)
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.runs[0]
        run.font.name = self.FONT_TITLE
        run.font.size = Pt(72)
        run.font.color.rgb = self.COLOR_BLUE  # 파란색
    
    def _add_header_box(self, slide, year, month):
        """헤더 박스 추가 (예산집행캘린더)"""
        # 파란색 세로 라인
        left_bar = slide.shapes.add_shape(
            1,  # 직사각형
            Inches(7.21), Inches(0.28),
            Inches(0.17), Inches(0.61)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = self.COLOR_BLUE
        left_bar.line.fill.background()
        
        # 헤더 박스
        header_box = slide.shapes.add_shape(
            1,  # 직사각형
            Inches(7.27), Inches(0.28),
            Inches(2.73), Inches(0.61)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = self.COLOR_WHITE
        header_box.line.fill.background()
        
        # 텍스트
        textbox = slide.shapes.add_textbox(
            Inches(7.49), Inches(0.39),
            Inches(2.73), Inches(0.40)
        )
        text_frame = textbox.text_frame
        text_frame.text = f"예산집행캘린더\n{str(year)[-2:]}.{month:02d}"
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = self.FONT_BODY
                run.font.size = Pt(16)
    
    def _add_calendar_table(self, slide, year, month, approval_days, execution_days, committee_day):
        """캘린더 표 생성"""
        # 월 캘린더 정보 가져오기
        cal = calendar.monthcalendar(year, month)
        rows = len(cal) + 1  # 헤더 포함
        cols = 7

        # 이전 달 마지막 날 계산
        if month == 1:
            prev_month = 12
            prev_year = year - 1
        else:
            prev_month = month - 1
            prev_year = year
        prev_month_last_day = calendar.monthrange(prev_year, prev_month)[1]

        # 표 추가
        table_shape = slide.shapes.add_table(
            rows, cols,
            self.TABLE_LEFT, self.TABLE_TOP,
            self.TABLE_WIDTH, self.TABLE_HEIGHT
        )
        table = table_shape.table

        # 표 스타일 비활성화 (기본 테두리가 적용되지 않도록)
        table.first_row = False
        table.first_col = False
        table.last_row = False
        table.last_col = False
        table.horz_banding = False
        table.vert_banding = False

        # 집행일, 결재일 위치 저장용
        execution_positions = []
        approval_positions = []

        # 행 높이 설정 (원본과 동일하게)
        table.rows[0].height = Inches(0.508)  # 헤더 행
        for row_idx in range(1, len(table.rows)):
            table.rows[row_idx].height = Inches(0.821)  # 날짜 행
        
        # 헤더 행 (요일)
        weekdays = ['Sun', 'Mon', 'Tue', 'Wed', 'Thu', 'Fri', 'Sat']
        for col_idx, day_name in enumerate(weekdays):
            cell = table.rows[0].cells[col_idx]
            cell.text = day_name
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = 1  # 중앙 정렬

            # 폰트 설정
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = self.FONT_BODY
            run.font.size = Pt(16)
            run.font.bold = True

            # 요일별 색상
            if col_idx == 0:  # 일요일
                run.font.color.rgb = self.COLOR_SUNDAY
            else:  # 평일 및 토요일
                run.font.color.rgb = self.COLOR_BLACK

            # 셀 배경 및 파란 테두리
            self._set_cell_background(cell, self.COLOR_WHITE)
            self._set_cell_border(cell, color_rgb=self._COLOR_BLUE_RGB, width=19050)  # 1.5pt
        
        # 날짜 행
        for row_idx, week in enumerate(cal, start=1):
            # 첫 번째 주에서 1일의 위치 찾기 (이전 달 날짜 계산용)
            if row_idx == 1:
                first_day_col = -1
                for i, d in enumerate(week):
                    if d == 1:
                        first_day_col = i
                        break

            for col_idx, day in enumerate(week):
                cell = table.rows[row_idx].cells[col_idx]

                # 지난달 날짜 계산
                if day == 0:
                    # 첫 번째 주에서 0인 경우 이전 달 날짜 표시
                    if row_idx == 1 and first_day_col > 0:
                        # 이전 달의 날짜를 역순으로 계산
                        prev_day = prev_month_last_day - (first_day_col - col_idx - 1)
                        cell.text = str(prev_day)
                        is_prev_month = True
                    else:
                        cell.text = ""
                        is_prev_month = False
                else:
                    cell.text = str(day)
                    is_prev_month = False

                # 셀 정렬
                if cell.text:
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = 1  # 중앙 정렬

                    # 폰트
                    run = cell.text_frame.paragraphs[0].runs[0]
                    run.font.name = self.FONT_BODY
                    run.font.size = Pt(20)

                    # 지난달 날짜
                    if is_prev_month:
                        run.font.color.rgb = self.COLOR_PREV_MONTH
                    # 집행일 표시
                    elif execution_days and day in execution_days:
                        run.font.color.rgb = self.COLOR_BLUE
                        run.font.bold = True
                        # 집행일 위치 저장 (row, col, day)
                        execution_positions.append((row_idx, col_idx, day))
                    else:
                        # 결재일 위치 저장 (토요일만 저장하여 금토 2칸 박스 표시)
                        # calendar.monthcalendar()에서 토요일은 col_idx=5
                        if approval_days and day in approval_days and col_idx == 5:  # 토요일
                            approval_positions.append((row_idx, col_idx, day))

                        # 요일별 색상
                        if col_idx == 0:  # 일요일
                            run.font.color.rgb = self.COLOR_SUNDAY
                        elif col_idx == 6:  # 토요일
                            run.font.color.rgb = self.COLOR_SATURDAY
                        else:  # 평일
                            run.font.color.rgb = self.COLOR_BLACK

                    # 운영위원회의 표시
                    if committee_day and day == committee_day and not is_prev_month:
                        cell.text = f"{day}\n운영위원회"
                        # 첫 줄 (날짜)
                        para1 = cell.text_frame.paragraphs[0]
                        para1.alignment = PP_ALIGN.CENTER
                        for run in para1.runs:
                            run.font.name = self.FONT_BODY
                            run.font.size = Pt(20)
                            if col_idx == 0:
                                run.font.color.rgb = self.COLOR_SUNDAY

                        # 두 번째 줄 (운영위원회) - 검정 글자, 황금색 강조(배경), 14pt
                        if len(cell.text_frame.paragraphs) > 1:
                            para2 = cell.text_frame.paragraphs[1]
                            para2.alignment = PP_ALIGN.CENTER
                            for run in para2.runs:
                                run.font.name = self.FONT_BODY
                                run.font.size = Pt(14)
                                run.font.color.rgb = self.COLOR_BLACK
                                run.font.highlight_color = self.COLOR_GOLD
                
                # 배경색 및 파란 테두리
                self._set_cell_background(cell, self.COLOR_WHITE)
                self._set_cell_border(cell, color_rgb=self._COLOR_BLUE_RGB, width=19050)  # 1.5pt

        return execution_positions, approval_positions

    def _set_cell_border(self, cell, color_rgb=(0, 0, 0), width=19050):
        """셀 테두리 설정 (width는 EMU 단위, 12700 = 1pt)"""
        from pptx.oxml.xmlchemy import OxmlElement

        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        # RGB 튜플을 16진수로 변환
        r, g, b = color_rgb
        hex_color = f"{r:02X}{g:02X}{b:02X}"

        # 테두리 설정 함수
        def add_border(border_name):
            # 기존 테두리 제거
            for child in list(tcPr):
                if child.tag.endswith(border_name):
                    tcPr.remove(child)

            # 새 테두리 추가
            ln = OxmlElement(f'a:{border_name}')
            ln.set('w', str(width))
            ln.set('cap', 'flat')
            ln.set('cmpd', 'sng')
            ln.set('algn', 'ctr')

            solidFill = OxmlElement('a:solidFill')
            srgbClr = OxmlElement('a:srgbClr')
            srgbClr.set('val', hex_color)
            solidFill.append(srgbClr)
            ln.append(solidFill)

            prstDash = OxmlElement('a:prstDash')
            prstDash.set('val', 'solid')
            ln.append(prstDash)

            tcPr.append(ln)

        # 모든 방향 테두리 추가
        add_border('lnL')  # 왼쪽
        add_border('lnR')  # 오른쪽
        add_border('lnT')  # 위
        add_border('lnB')  # 아래
    
    def _add_legend(self, slide, approval_days, execution_days, execution_positions, approval_positions):
        """범례 추가 - 집행일, 결재일 날짜 밑에 배치"""

        col_width = self.TABLE_WIDTH / 7  # 7개 열

        # 각 집행일 아래에 "예산집행일" 박스 추가
        if execution_positions:
            for row_idx, col_idx, day in execution_positions:
                # 셀의 X 위치 계산
                cell_left = self.TABLE_LEFT + (col_width * col_idx)
                cell_center_x = cell_left + (col_width / 2)

                # 셀의 Y 위치 계산 (표 시작 + 헤더 높이 + 날짜 행들 높이)
                # row_idx는 1부터 시작하므로 row_idx-1을 곱함
                cell_top = self.TABLE_TOP + Inches(0.508) + (Inches(0.821) * (row_idx - 1))
                cell_bottom = cell_top + Inches(0.821)

                # 박스 크기
                box_width = Inches(1.17)
                box_height = Inches(0.35)

                # 박스를 셀 아래에 배치 (살짝 내림)
                box_left = cell_center_x - (box_width / 2)
                box_top = cell_bottom - Inches(0.40)  # 살짝 내림

                # 예산집행일 박스 (둥근 모서리) - 텍스트 포함
                box = slide.shapes.add_shape(
                    5,  # 둥근 모서리 사각형
                    box_left, box_top,
                    box_width, box_height
                )
                # 모서리를 더 둥글게 설정 (0.0 ~ 0.5, 클수록 더 둥글음)
                box.adjustments[0] = 0.3
                box.fill.solid()
                box.fill.fore_color.rgb = self.COLOR_BLUE
                box.line.fill.background()

                # 도형 내부에 텍스트 추가
                text_frame = box.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = "예산집행일"
                p.alignment = PP_ALIGN.CENTER

                # 텍스트 수직 중앙 정렬
                text_frame.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE

                # 폰트 설정
                run = p.runs[0]
                run.font.name = self.FONT_BODY
                run.font.size = Pt(12)
                run.font.color.rgb = self.COLOR_WHITE
                run.font.bold = True

                # 예산집행일 옆에 설명 텍스트 추가
                desc_left = box_left + box_width + Inches(0.05)
                desc_top = box_top
                desc_width = Inches(4.0)
                desc_height = box_height

                desc_textbox = slide.shapes.add_textbox(
                    desc_left, desc_top,
                    desc_width, desc_height
                )
                desc_frame = desc_textbox.text_frame
                desc_frame.text = "<-전 주 토요일 자정까지 결재 난 건에 한해"
                desc_frame.vertical_anchor = 1  # 중앙 정렬

                desc_para = desc_frame.paragraphs[0]
                desc_para.alignment = PP_ALIGN.LEFT
                desc_run = desc_para.runs[0]
                desc_run.font.name = self.FONT_BODY
                desc_run.font.size = Pt(14)  # 운영위원회와 같은 크기
                desc_run.font.color.rgb = self.COLOR_BLACK
                desc_run.font.bold = True

        # 각 결재일 아래에 "결재일" 박스 추가 (금토 2칸에 걸침)
        if approval_positions:
            for row_idx, col_idx, day in approval_positions:
                # 금요일부터 토요일까지 2칸에 걸치도록 배치
                # calendar 데이터와 표 헤더가 불일치: calendar[Mon,Tue,Wed,Thu,Fri,Sat,Sun] vs 헤더[Sun,Mon,Tue,Wed,Thu,Fri,Sat]
                # 토요일 데이터(col_idx=5)는 헤더상 Fri 위치에 있으므로, Fri-Sat 박스를 그리려면 col_idx부터 시작
                friday_col_idx = col_idx  # 현재 위치(헤더상 Fri)부터 시작
                cell_left_friday = self.TABLE_LEFT + (col_width * friday_col_idx)

                # 셀의 Y 위치 계산
                cell_top = self.TABLE_TOP + Inches(0.508) + (Inches(0.821) * (row_idx - 1))
                cell_bottom = cell_top + Inches(0.821)

                # 박스 크기 (2칸 너비)
                box_width = col_width * 2
                box_height = Inches(0.35)

                # 박스를 금요일부터 토요일 2칸에 걸쳐 배치
                box_left = cell_left_friday
                box_top = cell_bottom - Inches(0.40)  # 살짝 내림

                # 결재일 박스 (둥근 모서리) - 텍스트 포함
                box = slide.shapes.add_shape(
                    5,  # 둥근 모서리 사각형
                    box_left, box_top,
                    box_width, box_height
                )
                # 모서리를 더 둥글게 설정 (0.0 ~ 0.5, 클수록 더 둥글음)
                box.adjustments[0] = 0.3
                box.fill.solid()
                box.fill.fore_color.rgb = self.COLOR_PINK
                box.line.fill.background()

                # 도형 내부에 텍스트 추가
                text_frame = box.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = "결재일"
                p.alignment = PP_ALIGN.CENTER

                # 텍스트 수직 중앙 정렬
                text_frame.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE

                # 폰트 설정
                run = p.runs[0]
                run.font.name = self.FONT_BODY
                run.font.size = Pt(12)
                run.font.color.rgb = self.COLOR_WHITE
                run.font.bold = True
        
        # 운영위원회의 주석
        footer = slide.shapes.add_textbox(
            Inches(0.64), Inches(6.80),
            Inches(8.0), Inches(0.34)
        )
        footer.text = "* 운영위원회의: 마지막주 전주 일요일"
        run = footer.text_frame.paragraphs[0].runs[0]
        run.font.name = self.FONT_BODY
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)
    
    def _format_legend_text(self, textbox, color):
        """범례 텍스트 포맷"""
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.runs[0]
        run.font.name = self.FONT_BODY
        run.font.size = Pt(14)
        run.font.color.rgb = color
    
    def _set_cell_background(self, cell, color):
        """셀 배경색 설정"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
    
    def save(self, filename):
        """PPT 저장"""
        self.prs.save(filename)


# ============================================
# 사용 예시
# ============================================

if __name__ == "__main__":
    # 템플릿 생성
    template = BudgetCalendarTemplate()
    
    # 방법 1: 자동 계산 (규정에 따라 자동으로 날짜 계산)
    print("=" * 60)
    print("자동 계산 모드 - 규정에 따른 일정 생성")
    print("=" * 60)
    
    for month in range(1, 13):
        schedule = BudgetCalendarTemplate.calculate_schedule(2026, month)
        print(f"\n{month}월:")
        print(f"  예산집행일 (둘째주/넷째주 화요일): {schedule['execution_days']}")
        print(f"  결재일 (전주 금요일/토요일): {schedule['approval_days']}")
        print(f"  운영위원회의 (마지막주 전주 일요일): {schedule['committee_day']}")
        
        template.create_calendar_slide(
            year=2026,
            month=month,
            auto_calculate=True  # 자동 계산
        )
    
    # 방법 2: 수동 지정 (특정 날짜를 직접 지정하고 싶을 때)
    # template.create_calendar_slide(
    #     year=2026,
    #     month=1,
    #     approval_days=[4, 18],
    #     execution_days=[6, 20],
    #     committee_day=26,
    #     auto_calculate=False
    # )
    
    # 저장
    import time
    timestamp = time.strftime("%Y%m%d_%H%M%S")
    output_file = f"2026_Budget_Calendar_{timestamp}.pptx"
    template.save(output_file)

    print(f"\n{'=' * 60}")
    print(f"PPT 생성 완료: {output_file}")
    print(f"슬라이드 수: {len(template.prs.slides)}")
    print(f"저장 위치: {os.path.abspath(output_file)}")
    print(f"{'=' * 60}")
